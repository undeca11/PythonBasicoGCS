from pptx import Presentation
from pptx.util import Inches

# Criar apresentação
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Slide de título
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Business Intelligence - IBM"
slide.placeholders[1].text = "Perguntas e Respostas com base no site da IBM"

# Perguntas e Respostas
qa_pairs = [
    ("1. O que é Business Intelligence (BI) e qual é seu objetivo principal?",
     "BI é o processo de coleta, armazenamento, análise e visualização de dados de negócios, com o objetivo de transformar dados brutos em informações acionáveis para embasar decisões estratégicas."),
    ("2. Quais são os principais componentes de uma solução de BI?",
     "Data warehouses, data lakes e data lakehouses. Esses elementos permitem armazenar, consolidar e analisar grandes volumes de dados de forma eficiente."),
    ("3. Como a IBM utiliza a arquitetura de malha de dados em BI?",
     "A IBM usa a malha de dados para orquestrar dados de forma inteligente e descentralizada, conectando os dados certos às pessoas certas no momento certo."),
    ("4. Quais são os benefícios do uso de BI nas organizações?",
     "Melhoria na tomada de decisões, identificação de tendências, otimização de processos e descoberta de novas oportunidades de negócios."),
    ("5. O que é o IBM Planning Analytics e como ele contribui para o BI?",
     "É uma solução baseada em IA que permite prever cenários e realizar análises dinâmicas, apoiando decisões com dados integrados em tempo real."),
    ("6. Como o IBM Cognos Analytics facilita a análise de dados?",
     "Oferece dashboards interativos, relatórios automatizados e análise preditiva, permitindo decisões baseadas em visualizações e dados confiáveis."),
    ("7. Qual é o papel da inteligência artificial (IA) no BI?",
     "Automatiza a análise, identifica padrões, melhora previsões e recomenda ações estratégicas, aumentando a precisão das decisões."),
    ("8. O que é BI generativa e como ela beneficia as empresas?",
     "Uso da IA para gerar insights e relatórios automaticamente. Acelera análises, reduz erros e entrega insights em tempo real."),
    ("9. Quais são os desafios enfrentados pelas organizações ao adotar BI?",
     "Integração de dados de diferentes fontes, manutenção da qualidade e capacitação dos usuários para uso eficaz das ferramentas."),
    ("10. Como a IBM apoia as organizações na implementação de BI?",
     "Com ferramentas como Cloud Pak for Data, Watson Studio e SPSS, além de consultorias que ajudam na coleta, análise e uso de dados.")
]

# Criar slides de conteúdo
for pergunta, resposta in qa_pairs:
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = pergunta
    slide.placeholders[1].text = resposta

# Salvar arquivo
prs.save("Business_Intelligence_IBM_Slides.pptx")
